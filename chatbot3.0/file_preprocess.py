from langchain_community.document_loaders import UnstructuredFileLoader
from typing import List
import tqdm
from langchain_community.document_loaders import PyPDFLoader
from docx import Document
import markdown2
from bs4 import BeautifulSoup
import pdfplumber
from pdf2image import convert_from_path
import requests
import html2text
import urllib.request

from typing import TYPE_CHECKING

if TYPE_CHECKING:
    try:
        from rapidocr_paddle import RapidOCR
    except ImportError:
        from rapidocr_onnxruntime import RapidOCR

def get_ocr(use_cuda: bool = True) -> "RapidOCR":
    try:
        from rapidocr_paddle import RapidOCR
        ocr = RapidOCR(det_use_cuda=use_cuda, cls_use_cuda=use_cuda, rec_use_cuda=use_cuda)
    except ImportError:
        from rapidocr_onnxruntime import RapidOCR
        ocr = RapidOCR()
    return ocr

def process_url(url):
    flag = False
    html_content = ''
    try:
        response = requests.get(url)
        html_content = response.text
    except:pass
    if len(html_content) < 1:
        try:
            response = urllib.request.urlopen(url)
            html_content = response.read().decode('utf-8')
        except:pass
    try:
        if len(html_content) > 0:
            html_content = html2text.html2text(html_content)
    except:pass
    html_content = html_content.strip()
    if len(html_content) > 0:
        flag = True
    return flag, html_content
# def process_pdf(pdf_path):  # v1.0与v2.0版本使用，v3.0已经弃用
#     loader = PyPDFLoader(pdf_path, extract_images=True)
#     pages = loader.load_and_split()
#     text = ''
#     for page in pages:
#         text+= page.page_content
#         text+='\n'
#     return text

def process_pdf(pdf_path):
    with pdfplumber.open(pdf_path) as pdf:
        all_text = ''  # 用于存储提取的文本内容
        for i,page in enumerate(pdf.pages):
            try:
                text = page.extract_text()
                if text:
                    lines = text.split('\n')
                    # 对于第二页，提取页眉内容
                    if i == 3:  # 页码从4开始，第二页的索引为1
                        header_text = lines[0]  # 假设页眉在第一行
                    # 从第一页开始去除页眉和页脚
                    text_without_header_footer = '\n'.join(lines[1:-1])  # 去除第一行和最后一行
                    all_text += text_without_header_footer + '\n\n'
                # 检查是否存在表格，并仅处理存在表格的页面
                if page.extract_tables():
                    markdown_content = ""
                    for table in page.extract_tables():
                        markdown_table = ''  # 存储当前表格的Markdown表示
                        for i, row in enumerate(table):
                            # 移除空列，这里假设空列完全为空，根据实际情况调整
                            row = [cell for cell in row if cell is not None and cell != '']
                            # 转换每个单元格内容为字符串，并用竖线分隔
                            processed_row = [str(cell).strip() if cell is not None else "" for cell in row]
                            markdown_row = '| ' + ' | '.join(processed_row) + ' |\n'
                            markdown_table += markdown_row
                            # 对于表头下的第一行，添加分隔线
                            if i == 0:
                                separators = [':---' if cell.isdigit() else '---' for cell in row]
                                markdown_table += '| ' + ' | '.join(separators) + ' |\n'
                        all_text += markdown_table + '\n'
            except Exception as e:
                # 进行ocr处理
                tmp_img =  convert_from_path(pdf_path, dpi=500, first_page=i+1, last_page=i+1) #获取第一页内容
                ocr = get_ocr()
                result, _ = ocr(tmp_img[0])
                ocr_text = ''
                if result:
                    ocr_result = [line[1] for line in result]
                    for res in ocr_result:
                        ocr_text+=res+'\n'
                all_text +=ocr_text+'\n'
                print("利用OCR处理第{}页文件".format(i)) 
    return all_text

def process_txt(txt_path):
    with open(txt_path,'r',encoding='utf-8') as file:
        text = file.read()
        return text
def process_word(word_path):
    # 打开Word文档
    doc = Document(word_path)
    text = ''
    for paragraph in doc.paragraphs:
        # 提取段落文本
        text_ = paragraph.text.strip()
        if text_:  # 忽略空段落
            text+=text_
            text+='\n'
    return text
def process_md(md_path):
    with open(md_path, 'r', encoding='utf-8') as f:
        content = f.read()
    # 将Markdown转换为HTML
    html_content = markdown2.markdown(content)
    soup = BeautifulSoup(html_content, "html.parser")
    text = soup.get_text()
    return text
class RapidOCRPPTLoader(UnstructuredFileLoader):
    def _get_elements(self) -> List:
        def ppt2text(filepath):
            from pptx import Presentation
            from PIL import Image
            import numpy as np
            from io import BytesIO
            from rapidocr_onnxruntime import RapidOCR
            ocr = RapidOCR()
            prs = Presentation(filepath)
            resp = ""

            def extract_text(shape):
                nonlocal resp
                if shape.has_text_frame:
                    resp += shape.text.strip() + "\n"
                if shape.has_table:
                    for row in shape.table.rows:
                        for cell in row.cells:
                            for paragraph in cell.text_frame.paragraphs:
                                resp += paragraph.text.strip() + "\n"
                if shape.shape_type == 13:  # 13 表示图片
                    image = Image.open(BytesIO(shape.image.blob))
                    result, _ = ocr(np.array(image))
                    if result:
                        ocr_result = [line[1] for line in result]
                        resp += "\n".join(ocr_result)
                elif shape.shape_type == 6:  # 6 表示组合
                    for child_shape in shape.shapes:
                        extract_text(child_shape)

            b_unit = tqdm.tqdm(total=len(prs.slides),
                               desc="RapidOCRPPTLoader slide index: 1")
            # 遍历所有幻灯片
            for slide_number, slide in enumerate(prs.slides, start=1):
                b_unit.set_description(
                    "RapidOCRPPTLoader slide index: {}".format(slide_number))
                b_unit.refresh()
                sorted_shapes = sorted(slide.shapes,
                                       key=lambda x: (x.top, x.left))  # 从上到下、从左到右遍历
                for shape in sorted_shapes:
                    extract_text(shape)
                b_unit.update(1)
            return resp

        text = ppt2text(self.file_path)
        from unstructured.partition.text import partition_text
        return partition_text(text=text, **self.unstructured_kwargs)

def process_pptx(pptx_path):
    loader = RapidOCRPPTLoader(file_path=pptx_path)
    pages = loader.load_and_split()
    text = ''
    for page in pages:
        text+= page.page_content
        text+='\n'
    return text
    
if __name__ =='__main__':
    print('开始测试')
    file = './files//恒立实业：2023年年度报告.pdf'
    # text1 = process_pdf1(file)
    # text2 = process_pdf(file)
    data = process_pptx('./files/复杂网络上SIVS传染病模型的脉冲最优控制_开题_2272009-牛东雷(一稿).pptx')
    print(data)
